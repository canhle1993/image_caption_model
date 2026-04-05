"""
Tạo file PowerPoint bảo vệ đồ án:
  SINH MÔ TẢ ẢNH TỰ ĐỘNG SỬ DỤNG HỌC SÂU
Chạy: /opt/anaconda3/envs/tf-metal/bin/python generate_slides.py
Output: BAO_VE_DO_AN.pptx
"""

import json, base64, io, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

os.chdir(os.path.dirname(os.path.abspath(__file__)))

# ─── Màu sắc ────────────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x1F, 0x38, 0x64)
MID_BLUE    = RGBColor(0x2E, 0x75, 0xB6)
LIGHT_BLUE  = RGBColor(0xBD, 0xD7, 0xEE)
ORANGE      = RGBColor(0xED, 0x7D, 0x31)
GREEN       = RGBColor(0x70, 0xAD, 0x47)
RED         = RGBColor(0xC0, 0x00, 0x00)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x00, 0x00, 0x00)
DARK_GREY   = RGBColor(0x26, 0x26, 0x26)
LIGHT_GREY  = RGBColor(0xF2, 0xF2, 0xF2)
YELLOW_HL   = RGBColor(0xFF, 0xFF, 0x00)

# ─── Kích thước slide 16:9 ──────────────────────────────────────────────────
W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK_LAYOUT = prs.slide_layouts[6]   # blank

# ─── Đọc ảnh từ notebook ────────────────────────────────────────────────────
def get_images(notebook_path):
    """Trả về list các (cell_index, img_index, base64_str, cell_title)"""
    nb = json.load(open(notebook_path, encoding='utf-8'))
    results = []
    for ci, cell in enumerate(nb['cells']):
        img_in_cell = 0
        for out in cell.get('outputs', []):
            data = out.get('data', {})
            if 'image/png' in data:
                src = cell.get('source', '')
                title = ''.join(src) if isinstance(src, list) else src
                lines = [l.strip() for l in title.split('\n') if l.strip()]
                label = lines[0] if lines else f'Cell {ci}'
                results.append((ci, img_in_cell, data['image/png'], label))
                img_in_cell += 1
    return results

imgs_vgg   = get_images('ImageCaptioning_Flickr30k_VGG16.ipynb')
imgs_res   = get_images('ImageCaptioning_Flickr30k_ResNet101.ipynb')
imgs_clip  = get_images('ImageCaptioning_Flickr30k_CLIP.ipynb')

def b64_to_stream(b64str):
    if isinstance(b64str, list):
        b64str = ''.join(b64str)
    return io.BytesIO(base64.b64decode(b64str))

# ─── Helpers ────────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill_rgb, transparency=0):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE.RECTANGLE
    shape.line.fill.background()
    shape.line.width = 0
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = fill_rgb
    if transparency:
        fill.fore_color.theme_color = None
        sp = shape._element
        solidFill = sp.find('.//' + qn('a:solidFill'))
        if solidFill is not None:
            srgb = solidFill.find(qn('a:srgbClr'))
            if srgb is not None:
                alpha = etree.SubElement(srgb, qn('a:alpha'))
                alpha.set('val', str(int((1 - transparency) * 100000)))
    return shape

def add_text(slide, text, x, y, w, h,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = 'Calibri'
    return txBox

def add_img(slide, b64str, x, y, w, h=None):
    stream = b64_to_stream(b64str)
    if h:
        slide.shapes.add_picture(stream, x, y, w, h)
    else:
        slide.shapes.add_picture(stream, x, y, w)

def title_bar(slide, title, subtitle=None):
    """Thanh tiêu đề xanh ở trên cùng"""
    add_rect(slide, 0, 0, W, Inches(1.1), DARK_BLUE)
    add_text(slide, title,
             Inches(0.25), Inches(0.08), W - Inches(0.5), Inches(0.7),
             font_size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.25), Inches(0.72), W - Inches(0.5), Inches(0.38),
                 font_size=14, bold=False, color=LIGHT_BLUE)

def footer(slide, text="Neural Image Caption Generation  |  DT2307L  |  Aptech"):
    add_rect(slide, 0, H - Inches(0.35), W, Inches(0.35), DARK_BLUE)
    add_text(slide, text, Inches(0.2), H - Inches(0.33), W - Inches(0.4), Inches(0.3),
             font_size=10, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

def slide_bg(slide):
    add_rect(slide, 0, 0, W, H, WHITE)

# ─── SLIDE 1: TRANG BÌA ─────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(sl, 0, 0, W, H, DARK_BLUE)
add_rect(sl, 0, Inches(1.5), W, Inches(0.06), ORANGE)  # gạch cam
add_rect(sl, 0, Inches(5.8), W, Inches(0.06), ORANGE)

add_text(sl, "SINH MÔ TẢ ẢNH TỰ ĐỘNG",
         Inches(0.5), Inches(1.7), W - Inches(1), Inches(1.0),
         font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "SỬ DỤNG HỌC SÂU",
         Inches(0.5), Inches(2.65), W - Inches(1), Inches(0.8),
         font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Neural Image Caption Generation",
         Inches(0.5), Inches(3.4), W - Inches(1), Inches(0.6),
         font_size=24, bold=False, color=ORANGE, align=PP_ALIGN.CENTER, italic=True)

info = [
    ("Sinh viên thực hiện", "Lê Văn Cảnh    |    Nguyễn Đức Trường Giang"),
    ("Lớp",                 "DT2307L"),
    ("Giảng viên hướng dẫn","Hồ Nhựt Minh"),
    ("Trường",              "Aptech"),
    ("Năm học",             "2025 – 2026"),
]
y0 = Inches(4.3)
for label, value in info:
    add_text(sl, f"{label}:", Inches(2.5), y0, Inches(3.2), Inches(0.35),
             font_size=14, bold=True, color=ORANGE, align=PP_ALIGN.RIGHT)
    add_text(sl, value, Inches(5.8), y0, Inches(5.5), Inches(0.35),
             font_size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT)
    y0 += Inches(0.35)

# ─── SLIDE 2: MỤC LỤC ───────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK_LAYOUT)
slide_bg(sl)
title_bar(sl, "NỘI DUNG TRÌNH BÀY")
footer(sl)

items = [
    ("01", "Đặt vấn đề", "Bài toán Image Captioning là gì? Ứng dụng thực tế"),
    ("02", "Dataset", "Flickr30k — 31,783 ảnh, 158,915 caption"),
    ("03", "Công nghệ sử dụng", "VGG16 / ResNet-101 / CLIP + LSTM + Bahdanau Attention"),
    ("04", "Kiến trúc mô hình", "Encoder-Decoder, Teacher Forcing, Beam Search"),
    ("05", "Quá trình huấn luyện", "Hyperparameters, Training curves 3 mô hình"),
    ("06", "Kết quả & Demo", "BLEU Score so sánh, ảnh thật + caption sinh ra"),
    ("07", "Phân tích & Kết luận", "Giải thích kết quả, hướng phát triển"),
]

x_num = Inches(0.5)
x_title = Inches(1.6)
x_sub = Inches(4.5)
y_start = Inches(1.35)
row_h = Inches(0.76)

for i, (num, title, sub) in enumerate(items):
    y = y_start + i * row_h
    bg_col = LIGHT_BLUE if i % 2 == 0 else WHITE
    add_rect(sl, Inches(0.3), y, W - Inches(0.6), row_h - Inches(0.05), bg_col)
    add_text(sl, num, x_num, y + Inches(0.12), Inches(0.9), Inches(0.5),
             font_size=22, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    add_text(sl, title, x_title, y + Inches(0.05), Inches(2.8), Inches(0.4),
             font_size=17, bold=True, color=DARK_BLUE)
    add_text(sl, sub, x_title, y + Inches(0.38), Inches(8.5), Inches(0.35),
             font_size=12, color=DARK_GREY)

# ─── SLIDE 3: ĐẶT VẤN ĐỀ ────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK_LAYOUT)
slide_bg(sl)
title_bar(sl, "ĐẶT VẤN ĐỀ", "Image Captioning — Giao thoa giữa Computer Vision và NLP")
footer(sl)

add_text(sl, "Máy tính có thể nhìn ảnh và tự mô tả bằng ngôn ngữ tự nhiên không?",
         Inches(0.4), Inches(1.25), W - Inches(0.8), Inches(0.5),
         font_size=18, bold=True, color=DARK_BLUE)

add_rect(sl, Inches(0.4), Inches(1.85), Inches(5.8), Inches(2.3), LIGHT_BLUE)
add_text(sl, "🖼️  Ảnh đầu vào",
         Inches(0.6), Inches(1.9), Inches(5.4), Inches(0.4),
         font_size=15, bold=True, color=DARK_BLUE)
add_text(sl, "[ Một người đang đạp xe trên đường ]",
         Inches(0.6), Inches(2.25), Inches(5.4), Inches(0.5),
         font_size=14, color=DARK_GREY, italic=True)
add_text(sl, "⬇️  Mô hình AI",
         Inches(0.6), Inches(2.75), Inches(5.4), Inches(0.4),
         font_size=15, bold=True, color=DARK_BLUE)
add_text(sl, '"A man is riding a bicycle on the street"',
         Inches(0.6), Inches(3.1), Inches(5.4), Inches(0.5),
         font_size=16, bold=True, color=GREEN)

add_text(sl, "Ứng dụng thực tế:", Inches(6.5), Inches(1.85), Inches(6.4), Inches(0.4),
         font_size=16, bold=True, color=DARK_BLUE)
apps = [
    ("👁️", "Hỗ trợ người khiếm thị", "Đọc ảnh qua giọng nói"),
    ("🔍", "Tìm kiếm ảnh", "Tìm bằng mô tả ngôn ngữ tự nhiên"),
    ("📱", "Mạng xã hội", "Gợi ý caption cho ảnh đăng tải"),
    ("🤖", "Robot thông minh", "Hiểu và mô tả môi trường"),
    ("🏥", "Y tế", "Mô tả X-quang, MRI tự động"),
]
for j, (icon, title_a, sub) in enumerate(apps):
    y = Inches(2.3) + j * Inches(0.65)
    add_text(sl, f"{icon}  {title_a}", Inches(6.6), y, Inches(3.5), Inches(0.35),
             font_size=14, bold=True, color=DARK_BLUE)
    add_text(sl, sub, Inches(6.6), y + Inches(0.3), Inches(6.0), Inches(0.3),
             font_size=12, color=DARK_GREY)

add_rect(sl, Inches(0.4), Inches(4.3), W - Inches(0.8), Inches(0.6), DARK_BLUE)
add_text(sl, "Đây là bài toán kết hợp  Computer Vision  +  Natural Language Processing",
         Inches(0.6), Inches(4.37), W - Inches(1.2), Inches(0.45),
         font_size=16, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

# ─── SLIDE 4: MỤC TIÊU ──────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK_LAYOUT)
slide_bg(sl)
title_bar(sl, "MỤC TIÊU ĐỀ TÀI", "So sánh 3 hướng tiếp cận — Cùng decoder, cùng dataset, cùng tập test")
footer(sl)

models = [
    ("MÔ HÌNH 1", "VGG16 + LSTM", "Bahdanau Attention",
     "• Encoder: VGG16 (conv5_3)\n• Feature: (49, 512)\n• Vai trò: Baseline cơ sở\n• Paper: Simonyan & Zisserman, 2014", MID_BLUE),
    ("MÔ HÌNH 2", "ResNet-101 + LSTM", "Bahdanau Attention",
     "• Encoder: ResNet-101 (layer4)\n• Feature: (49, 2048)\n• Ưu điểm: Residual connections\n• Paper: He et al., 2016", GREEN),
    ("MÔ HÌNH 3", "CLIP ViT-B/32 + LSTM", "Bahdanau Attention",
     "• Encoder: CLIP Vision Transformer\n• Feature: (49, 512)\n• Huấn luyện trên 400M cặp ảnh-văn bản\n• Paper: Radford et al., 2021 (OpenAI)", ORANGE),
]
col_w = Inches(3.9)
for j, (num, title_m, sub_m, detail, color) in enumerate(models):
    x = Inches(0.4) + j * (col_w + Inches(0.25))
    add_rect(sl, x, Inches(1.25), col_w, Inches(0.5), color)
    add_text(sl, num, x, Inches(1.3), col_w, Inches(0.4),
             font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, x, Inches(1.75), col_w, Inches(0.55), DARK_BLUE)
    add_text(sl, title_m, x, Inches(1.8), col_w, Inches(0.45),
             font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, sub_m, x, Inches(2.3), col_w, Inches(0.35),
             font_size=13, color=color, align=PP_ALIGN.CENTER, bold=True)
    add_rect(sl, x, Inches(2.65), col_w, Inches(2.2), LIGHT_GREY)
    add_text(sl, detail, x + Inches(0.1), Inches(2.75), col_w - Inches(0.2), Inches(2.0),
             font_size=13, color=DARK_GREY)

add_rect(sl, Inches(0.4), Inches(5.05), W - Inches(0.8), Inches(0.6), DARK_BLUE)
add_text(sl, "✅  Cả 3 mô hình: cùng Dataset Flickr30k  |  cùng Decoder  |  cùng tập Test 1,000 ảnh  →  SO SÁNH CÔNG BẰNG",
         Inches(0.6), Inches(5.13), W - Inches(1.2), Inches(0.45),
         font_size=15, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

# ─── SLIDE 5: DATASET ────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK_LAYOUT)
slide_bg(sl)
title_bar(sl, "DATASET: FLICKR30K", "Young et al., 2014  —  Dataset chuẩn cho bài toán Image Captioning")
footer(sl)

# Bảng thống kê
headers = ["Thông số", "Giá trị"]
rows = [
    ["Tổng số ảnh", "31,783 ảnh"],
    ["Số caption", "158,915 câu  (5 câu / ảnh)"],
    ["Định dạng", "JPEG — đa dạng chủ đề"],
    ["Nguồn nhãn", "Amazon Mechanical Turk (người thật)"],
    ["Ngôn ngữ", "Tiếng Anh"],
]
split_rows = [
    ["Train", "29,769 ảnh", "93.7%", "148,844 caption"],
    ["Validation", "1,014 ảnh", "3.2%", "5,070 caption"],
    ["Test", "1,000 ảnh", "3.1%", "5,000 caption"],
]

xt = Inches(0.4)
yt = Inches(1.3)
col_ws2 = [Inches(3.2), Inches(4.2)]
add_rect(sl, xt, yt, sum(col_ws2), Inches(0.4), DARK_BLUE)
for j2, h in enumerate(headers):
    add_text(sl, h, xt + sum(col_ws2[:j2]), yt + Inches(0.05),
             col_ws2[j2], Inches(0.3), font_size=13, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
for ri, row in enumerate(rows):
    yy = yt + Inches(0.4) + ri * Inches(0.38)
    bg = LIGHT_GREY if ri % 2 == 0 else WHITE
    add_rect(sl, xt, yy, sum(col_ws2), Inches(0.38), bg)
    for j2, val in enumerate(row):
        add_text(sl, val, xt + sum(col_ws2[:j2]) + Inches(0.05), yy + Inches(0.05),
                 col_ws2[j2] - Inches(0.1), Inches(0.3),
                 font_size=12, color=DARK_GREY,
                 align=PP_ALIGN.CENTER if j2 == 1 else PP_ALIGN.LEFT)

# Phân chia
yt2 = Inches(3.7)
add_text(sl, "Phân chia dữ liệu (seed = 42):", xt, yt2, Inches(7.5), Inches(0.35),
         font_size=14, bold=True, color=DARK_BLUE)
split_headers = ["Tập", "Số ảnh", "Tỷ lệ", "Số caption"]
split_col_ws = [Inches(1.5), Inches(1.5), Inches(1.2), Inches(3.7)]
yt2 += Inches(0.35)
add_rect(sl, xt, yt2, sum(split_col_ws), Inches(0.38), DARK_BLUE)
for j2, h in enumerate(split_headers):
    add_text(sl, h, xt + sum(split_col_ws[:j2]), yt2 + Inches(0.05),
             split_col_ws[j2], Inches(0.28), font_size=12, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
for ri, row in enumerate(split_rows):
    yy = yt2 + Inches(0.38) + ri * Inches(0.38)
    bg_c = [LIGHT_BLUE, LIGHT_GREY, LIGHT_GREY][ri]
    add_rect(sl, xt, yy, sum(split_col_ws), Inches(0.38), bg_c)
    for j2, val in enumerate(row):
        add_text(sl, val, xt + sum(split_col_ws[:j2]) + Inches(0.05), yy + Inches(0.05),
                 split_col_ws[j2] - Inches(0.1), Inches(0.28),
                 font_size=12, color=DARK_GREY, bold=(ri == 0),
                 align=PP_ALIGN.CENTER)

# Ảnh mẫu từ notebook
if imgs_vgg and len(imgs_vgg) > 0:
    add_text(sl, "Ảnh mẫu từ dataset:", Inches(8.0), Inches(1.3), Inches(5.0), Inches(0.35),
             font_size=14, bold=True, color=DARK_BLUE)
    add_img(sl, imgs_vgg[0][2], Inches(8.0), Inches(1.65), Inches(5.1), Inches(3.2))

# ─── SLIDE 6: EDA — PHÂN TÍCH ẢNH ─────────────────────────────────────────
sl = prs.slides.add_slide(BLANK_LAYOUT)
slide_bg(sl)
title_bar(sl, "PHÂN TÍCH DỮ LIỆU (EDA)", "Phân bố kích thước ảnh và đặc điểm dataset")
footer(sl)

if len(imgs_vgg) > 1:
    add_text(sl, "Phân bố kích thước ảnh trong Flickr30k:", Inches(0.4), Inches(1.2),
             W - Inches(0.8), Inches(0.4), font_size=15, bold=True, color=DARK_BLUE)
    add_img(sl, imgs_vgg[1][2], Inches(0.3), Inches(1.6), W - Inches(0.6), Inches(3.8))
    add_text(sl, "→ Kích thước ảnh rất đa dạng (không đồng nhất), trung bình 459×393 px. Resize về 224×224 trước khi đưa vào encoder.",
             Inches(0.4), Inches(5.5), W - Inches(0.8), Inches(0.5),
             font_size=13, color=DARK_GREY, italic=True)

# ─── SLIDE 7: EDA — PHÂN TÍCH CAPTION ──────────────────────────────────────
sl = prs.slides.add_slide(BLANK_LAYOUT)
slide_bg(sl)
title_bar(sl, "PHÂN TÍCH DỮ LIỆU (EDA)", "Thống kê caption và từ vựng")
footer(sl)

if len(imgs_vgg) > 4:
    add_img(sl, imgs_vgg[4][2], Inches(0.3), Inches(1.15), W - Inches(0.6), Inches(4.2))
    add_text(sl, "Trước lọc: 23,457 từ  →  Sau lọc: 10,000 từ phổ biến nhất  |  Độ dài caption trung bình: ~12 từ",
             Inches(0.4), Inches(5.45), W - Inches(0.8), Inches(0.45),
             font_size=13, color=DARK_GREY, bold=True)

# ─── SLIDE 8: TIỀN XỬ LÝ CAPTION ───────────────────────────────────────────
sl = prs.slides.add_slide(BLANK_LAYOUT)
slide_bg(sl)
title_bar(sl, "TIỀN XỬ LÝ DỮ LIỆU", "Pipeline làm sạch caption + GloVe Embeddings")
footer(sl)

add_text(sl, "Quy trình làm sạch caption:", Inches(0.4), Inches(1.2), Inches(6), Inches(0.4),
         font_size=15, bold=True, color=DARK_BLUE)
steps = [
    ('1', 'Chuyển thành chữ thường (lowercase)',   '"Running" → "running"'),
    ('2', 'Loại bỏ ký tự đặc biệt',                '"park!!" → "park"'),
    ('3', 'Giữ từ có độ dài ≥ 1 ký tự',           'giữ "a", "i" — không loại'),
    ('4', 'Thêm token đặc biệt',                   '"startseq ... endseq"'),
    ('5', 'Lưu vào descriptions_flickr30k.txt',    '29,769 × 5 = 148,844 câu'),
]
for j, (num, step, ex) in enumerate(steps):
    y = Inches(1.7) + j * Inches(0.52)
    add_rect(sl, Inches(0.4), y, Inches(0.45), Inches(0.42), DARK_BLUE)
    add_text(sl, num, Inches(0.4), y + Inches(0.05), Inches(0.45), Inches(0.32),
             font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, step, Inches(0.95), y + Inches(0.05), Inches(4.5), Inches(0.32),
             font_size=13, bold=True, color=DARK_GREY)
    add_text(sl, ex, Inches(5.5), y + Inches(0.05), Inches(3.5), Inches(0.32),
             font_size=12, color=MID_BLUE, italic=True)

add_text(sl, "GloVe 300d  (Pennington et al., 2014):", Inches(0.4), Inches(4.55), Inches(9), Inches(0.4),
         font_size=15, bold=True, color=DARK_BLUE)
glove_stats = [
    ("Huấn luyện trên", "6 tỷ token từ Wikipedia + Gigaword"),
    ("Vector dimension", "300 chiều / từ"),
    ("Vocab coverage", "99.2%  từ vựng có sẵn GloVe vector"),
    ("Embedding matrix", "(10,004 × 300)  —  fine-tunable"),
]
for j, (k, v) in enumerate(glove_stats):
    y = Inches(5.0) + j * Inches(0.36)
    add_text(sl, f"• {k}:", Inches(0.6), y, Inches(2.8), Inches(0.32),
             font_size=13, bold=True, color=DARK_GREY)
    add_text(sl, v, Inches(3.5), y, Inches(5.5), Inches(0.32),
             font_size=13, color=DARK_GREY)

# Ảnh GloVe / tokenization
if len(imgs_vgg) > 5:
    add_img(sl, imgs_vgg[5][2], Inches(9.2), Inches(1.2), Inches(3.9), Inches(4.8))

# ─── SLIDE 9: KIẾN TRÚC TỔNG QUAN ───────────────────────────────────────────
sl = prs.slides.add_slide(BLANK_LAYOUT)
slide_bg(sl)
title_bar(sl, "KIẾN TRÚC TỔNG QUAN", "Encoder – Decoder với Bahdanau Attention")
footer(sl)

if len(imgs_vgg) > 7:
    add_text(sl, "Sơ đồ kiến trúc CaptionDecoder:", Inches(0.4), Inches(1.2), Inches(6), Inches(0.4),
             font_size=15, bold=True, color=DARK_BLUE)
    add_img(sl, imgs_vgg[7][2], Inches(0.3), Inches(1.6), Inches(7.5), Inches(4.3))

right_x = Inches(8.0)
add_text(sl, "3 thành phần chính:", right_x, Inches(1.2), Inches(5), Inches(0.4),
         font_size=14, bold=True, color=DARK_BLUE)
components = [
    (ORANGE,    "ENCODER (frozen)",   "CNN / ViT trích xuất feature (49, D)"),
    (MID_BLUE,  "DECODER (trainable)", "LSTMCell 512 + GloVe Embedding 300d"),
    (GREEN,     "ATTENTION",           "Bahdanau — tập trung đúng vùng ảnh"),
]
for j, (col, name, desc) in enumerate(components):
    y = Inches(1.7) + j * Inches(0.9)
    add_rect(sl, right_x, y, Inches(5.0), Inches(0.85), col)
    add_text(sl, name, right_x + Inches(0.1), y + Inches(0.05), Inches(4.8), Inches(0.35),
             font_size=14, bold=True, color=WHITE)
    add_text(sl, desc, right_x + Inches(0.1), y + Inches(0.4), Inches(4.8), Inches(0.35),
             font_size=12, color=WHITE)

add_rect(sl, right_x, Inches(4.65), Inches(5.0), Inches(0.55), DARK_BLUE)
add_text(sl, "~11.9M tham số  |  Loss: CrossEntropy + label_smoothing=0.1",
         right_x + Inches(0.1), Inches(4.72), Inches(4.8), Inches(0.4),
         font_size=12, color=ORANGE, bold=True)

# ─── SLIDE 10: BAHDANAU ATTENTION ───────────────────────────────────────────
sl = prs.slides.add_slide(BLANK_LAYOUT)
slide_bg(sl)
title_bar(sl, "CƠ CHẾ BAHDANAU ATTENTION", "Bahdanau et al., 2014  —  Trích dẫn >20,000 lần")
footer(sl)

add_text(sl, "Tại sao cần Attention?", Inches(0.4), Inches(1.2), Inches(6), Inches(0.4),
         font_size=15, bold=True, color=DARK_BLUE)
add_rect(sl, Inches(0.4), Inches(1.65), Inches(5.8), Inches(1.0), LIGHT_GREY)
add_text(sl, "Không có Attention: Toàn bộ ảnh → 1 vector cố định → LSTM chỉ 'nhớ' tổng quát",
         Inches(0.5), Inches(1.7), Inches(5.6), Inches(0.38),
         font_size=13, color=RED, bold=True)
add_text(sl, "Có Attention: Mỗi bước sinh từ → nhìn vào đúng vùng ảnh liên quan",
         Inches(0.5), Inches(2.1), Inches(5.6), Inches(0.38),
         font_size=13, color=GREEN, bold=True)

add_text(sl, "Công thức Bahdanau Attention:", Inches(0.4), Inches(2.85), Inches(6), Inches(0.4),
         font_size=15, bold=True, color=DARK_BLUE)
add_rect(sl, Inches(0.4), Inches(3.3), Inches(6.0), Inches(1.9), RGBColor(0xE8, 0xF4, 0xFF))
formulas = [
    "e_i   =   V^T · tanh( W_h · h_{t-1}  +  W_f · f_i )",
    "α_i   =   softmax( e_i )                              ← trọng số chú ý",
    "c_t   =   Σ α_i · f_i                                  ← context vector",
]
for j, f in enumerate(formulas):
    add_text(sl, f, Inches(0.6), Inches(3.4) + j * Inches(0.55), Inches(5.6), Inches(0.45),
             font_size=13, color=DARK_BLUE, bold=(j == 2))

legends = [
    ("h_{t-1}", "Trạng thái ẩn LSTM bước trước"),
    ("f_i", "Feature ảnh tại vị trí không gian i  (i = 1..49)"),
    ("α_i", "Mức độ chú ý tại vị trí i  (tổng α = 1.0)"),
    ("c_t", "Context vector — thông tin ảnh cần thiết tại bước t"),
]
for j, (sym, desc) in enumerate(legends):
    y = Inches(5.3) + j * Inches(0.32)
    add_text(sl, sym, Inches(0.5), y, Inches(1.0), Inches(0.28),
             font_size=12, bold=True, color=DARK_BLUE)
    add_text(sl, f":  {desc}", Inches(1.5), y, Inches(5.0), Inches(0.28),
             font_size=12, color=DARK_GREY)

# Ảnh Teacher Forcing structure
if len(imgs_vgg) > 6:
    add_text(sl, "Teacher Forcing — cấu trúc training sample:", Inches(6.5), Inches(1.2), Inches(6.5), Inches(0.4),
             font_size=14, bold=True, color=DARK_BLUE)
    add_img(sl, imgs_vgg[6][2], Inches(6.4), Inches(1.6), Inches(6.7), Inches(3.5))

# ─── SLIDE 11: VGG16 ENCODER ────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK_LAYOUT)
slide_bg(sl)
title_bar(sl, "ENCODER 1: VGG16", "Simonyan & Zisserman, 2014  —  Giải nhất ImageNet 2014")
footer(sl)

vgg_specs = [
    ("Kiến trúc",      "16 lớp CNN — đơn giản, đồng nhất"),
    ("Lớp trích xuất", "conv5_3  (block Conv cuối)"),
    ("Feature shape",  "(49, 512)  =  7×7 vị trí × 512 kênh"),
    ("Pre-trained",    "ImageNet — 1.2M ảnh, 1,000 lớp"),
    ("Encoder frozen", "Không fine-tune — tiết kiệm tài nguyên"),
    ("Vai trò",        "BASELINE — mô hình cơ sở để so sánh"),
]
for j, (k, v) in enumerate(vgg_specs):
    y = Inches(1.25) + j * Inches(0.52)
    add_rect(sl, Inches(0.4), y, Inches(2.0), Inches(0.44), DARK_BLUE)
    add_text(sl, k, Inches(0.45), y + Inches(0.07), Inches(1.9), Inches(0.3),
             font_size=12, bold=True, color=WHITE)
    add_text(sl, v, Inches(2.5), y + Inches(0.07), Inches(5.2), Inches(0.3),
             font_size=13, color=DARK_GREY)

add_text(sl, "Ảnh  224×224  →  [5 Conv Blocks]  →  Feature Map  7×7×512  →  reshape  →  (49, 512)",
         Inches(0.4), Inches(4.45), Inches(7.4), Inches(0.45),
         font_size=13, bold=True, color=MID_BLUE)

# VGG16 feature visualization
if len(imgs_vgg) > 3:
    add_text(sl, "Feature Map Visualization (VGG16 conv5_3):",
             Inches(7.8), Inches(1.2), Inches(5.3), Inches(0.38),
             font_size=14, bold=True, color=DARK_BLUE)
    add_img(sl, imgs_vgg[3][2], Inches(7.8), Inches(1.6), Inches(5.3), Inches(3.6))

# ─── SLIDE 12: RESNET-101 ENCODER ───────────────────────────────────────────
sl = prs.slides.add_slide(BLANK_LAYOUT)
slide_bg(sl)
title_bar(sl, "ENCODER 2: RESNET-101", "He et al., 2016  —  Giải nhất ImageNet 2015, CVPR Best Paper")
footer(sl)

res_specs = [
    ("Kiến trúc",      "101 lớp — Residual Connections (skip connection)"),
    ("Lớp trích xuất", "layer4  (block cuối)"),
    ("Feature shape",  "(49, 2048)  —  gấp 4 lần VGG16 (512d)"),
    ("Pre-trained",    "ImageNet — 1.2M ảnh"),
    ("Encoder frozen", "Không fine-tune"),
    ("Ưu điểm chính",  "Giải quyết vanishing gradient — học sâu không mất thông tin"),
]
for j, (k, v) in enumerate(res_specs):
    y = Inches(1.25) + j * Inches(0.52)
    add_rect(sl, Inches(0.4), y, Inches(2.2), Inches(0.44), GREEN)
    add_text(sl, k, Inches(0.45), y + Inches(0.07), Inches(2.1), Inches(0.3),
             font_size=12, bold=True, color=WHITE)
    add_text(sl, v, Inches(2.7), y + Inches(0.07), Inches(5.0), Inches(0.3),
             font_size=13, color=DARK_GREY)

add_rect(sl, Inches(0.4), Inches(4.45), Inches(7.4), Inches(0.55), LIGHT_GREY)
add_text(sl, "VGG16: (49, 512)  ←  ít thông tin    |    ResNet-101: (49, 2048)  ←  phong phú gấp 4",
         Inches(0.5), Inches(4.52), Inches(7.2), Inches(0.4),
         font_size=13, bold=True, color=DARK_BLUE)

if len(imgs_res) > 2:
    add_text(sl, "Feature Map Visualization (ResNet-101 layer4):",
             Inches(7.8), Inches(1.2), Inches(5.3), Inches(0.38),
             font_size=14, bold=True, color=DARK_BLUE)
    add_img(sl, imgs_res[2][2], Inches(7.8), Inches(1.6), Inches(5.3), Inches(3.6))

# ─── SLIDE 13: CLIP ENCODER ─────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK_LAYOUT)
slide_bg(sl)
title_bar(sl, "ENCODER 3: CLIP ViT-B/32", "Radford et al., 2021 (OpenAI)  —  Contrastive Language-Image Pretraining")
footer(sl)

clip_specs = [
    ("Kiến trúc",       "Vision Transformer (ViT-B/32) — patch 32×32"),
    ("Feature shape",   "(49, 512)"),
    ("Huấn luyện trên", "400 TRIỆU cặp (ảnh, văn bản) từ internet"),
    ("Phương pháp",     "Contrastive Learning — ảnh và caption 'gần nhau' trong không gian vector"),
    ("Đặc điểm đặc biệt","Feature đã chứa ngữ nghĩa NGÔN NGỮ — không chỉ visual"),
    ("Encoder frozen",  "Không fine-tune"),
]
for j, (k, v) in enumerate(clip_specs):
    y = Inches(1.25) + j * Inches(0.52)
    add_rect(sl, Inches(0.4), y, Inches(2.3), Inches(0.44), ORANGE)
    add_text(sl, k, Inches(0.45), y + Inches(0.07), Inches(2.2), Inches(0.3),
             font_size=12, bold=True, color=WHITE)
    add_text(sl, v, Inches(2.8), y + Inches(0.07), Inches(4.9), Inches(0.3),
             font_size=12, color=DARK_GREY)

add_rect(sl, Inches(0.4), Inches(4.45), Inches(7.4), Inches(0.65), DARK_BLUE)
add_text(sl, "VGG16 / ResNet: Chỉ học từ ảnh (ImageNet)\nCLIP: Học đồng thời từ ẢNH + VĂN BẢN → Feature mang ngữ nghĩa ngôn ngữ",
         Inches(0.5), Inches(4.48), Inches(7.2), Inches(0.6),
         font_size=13, bold=True, color=ORANGE)

if len(imgs_clip) > 1:
    add_text(sl, "Feature Map Visualization (CLIP ViT-B/32):",
             Inches(7.8), Inches(1.2), Inches(5.3), Inches(0.38),
             font_size=14, bold=True, color=DARK_BLUE)
    add_img(sl, imgs_clip[1][2], Inches(7.8), Inches(1.6), Inches(5.3), Inches(3.6))

# ─── SLIDE 14: TRAINING VGG16 ────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK_LAYOUT)
slide_bg(sl)
title_bar(sl, "HUẤN LUYỆN — MÔ HÌNH 1: VGG16", "30 epochs  |  Optimizer: AdamW  |  LR: 1e-4")
footer(sl)

if len(imgs_vgg) > 8:
    add_img(sl, imgs_vgg[8][2], Inches(0.3), Inches(1.15), Inches(9.0), Inches(4.3))

right_x2 = Inches(9.5)
add_text(sl, "Kết quả huấn luyện:", right_x2, Inches(1.2), Inches(3.6), Inches(0.4),
         font_size=14, bold=True, color=DARK_BLUE)
vgg_train = [
    ("Epochs thực tế", "30"),
    ("Best Val Loss", "3.8851"),
    ("Best Epoch", "23"),
    ("LR cuối cùng", "6.25e-6"),
    ("Batch size", "64"),
    ("Early stopping", "patience = 7"),
]
for j, (k, v) in enumerate(vgg_train):
    y = Inches(1.65) + j * Inches(0.48)
    add_rect(sl, right_x2, y, Inches(3.6), Inches(0.42), LIGHT_GREY if j%2==0 else WHITE)
    add_text(sl, k, right_x2 + Inches(0.1), y + Inches(0.07), Inches(2.0), Inches(0.28),
             font_size=12, bold=True, color=DARK_BLUE)
    add_text(sl, v, right_x2 + Inches(2.2), y + Inches(0.07), Inches(1.3), Inches(0.28),
             font_size=12, bold=True, color=MID_BLUE, align=PP_ALIGN.RIGHT)

# ─── SLIDE 15: TRAINING RESNET-101 ──────────────────────────────────────────
sl = prs.slides.add_slide(BLANK_LAYOUT)
slide_bg(sl)
title_bar(sl, "HUẤN LUYỆN — MÔ HÌNH 2: RESNET-101", "10 epochs  |  Optimizer: AdamW  |  LR: 1e-4")
footer(sl)

if len(imgs_res) > 4:
    add_img(sl, imgs_res[4][2], Inches(0.3), Inches(1.15), Inches(9.0), Inches(4.3))

right_x2 = Inches(9.5)
add_text(sl, "Kết quả huấn luyện:", right_x2, Inches(1.2), Inches(3.6), Inches(0.4),
         font_size=14, bold=True, color=DARK_BLUE)
res_train = [
    ("Epochs thực tế", "10"),
    ("Best Val Loss", "3.8775"),
    ("Thời gian/epoch", "~1,200 s"),
    ("Feature extraction", "8 img/s"),
    ("Batch size", "64"),
    ("Hội tụ", "Nhanh hơn VGG16"),
]
for j, (k, v) in enumerate(res_train):
    y = Inches(1.65) + j * Inches(0.48)
    add_rect(sl, right_x2, y, Inches(3.6), Inches(0.42), LIGHT_GREY if j%2==0 else WHITE)
    add_text(sl, k, right_x2 + Inches(0.1), y + Inches(0.07), Inches(2.0), Inches(0.28),
             font_size=12, bold=True, color=DARK_BLUE)
    add_text(sl, v, right_x2 + Inches(2.2), y + Inches(0.07), Inches(1.3), Inches(0.28),
             font_size=12, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)

# ─── SLIDE 16: TRAINING CLIP ─────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK_LAYOUT)
slide_bg(sl)
title_bar(sl, "HUẤN LUYỆN — MÔ HÌNH 3: CLIP", "10 epochs  |  Optimizer: AdamW  |  LR: 1e-4")
footer(sl)

if len(imgs_clip) > 4:
    add_img(sl, imgs_clip[4][2], Inches(0.3), Inches(1.15), Inches(9.0), Inches(4.3))

right_x2 = Inches(9.5)
add_text(sl, "Kết quả huấn luyện:", right_x2, Inches(1.2), Inches(3.6), Inches(0.4),
         font_size=14, bold=True, color=DARK_BLUE)
clip_train = [
    ("Epochs thực tế", "10"),
    ("Best Val Loss", "3.8817"),
    ("Val Accuracy", "40.1%  ← cao nhất"),
    ("Feature extraction", "122 img/s  ← nhanh nhất"),
    ("Batch size", "64"),
    ("Hội tụ", "Nhanh, smooth"),
]
for j, (k, v) in enumerate(clip_train):
    y = Inches(1.65) + j * Inches(0.48)
    add_rect(sl, right_x2, y, Inches(3.6), Inches(0.42), LIGHT_GREY if j%2==0 else WHITE)
    add_text(sl, k, right_x2 + Inches(0.1), y + Inches(0.07), Inches(2.0), Inches(0.28),
             font_size=12, bold=True, color=DARK_BLUE)
    add_text(sl, v, right_x2 + Inches(2.2), y + Inches(0.07), Inches(1.3), Inches(0.28),
             font_size=12, bold=True, color=ORANGE, align=PP_ALIGN.RIGHT)

# ─── SLIDE 17: BLEU SCORE ────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK_LAYOUT)
slide_bg(sl)
title_bar(sl, "THƯỚC ĐO BLEU SCORE", "Papineni et al., 2002 (IBM)  —  Chuẩn đánh giá trong hơn 20 năm")
footer(sl)

add_text(sl, "BLEU so sánh n-gram của caption sinh ra với ground truth:", Inches(0.4), Inches(1.2),
         W - Inches(0.8), Inches(0.4), font_size=15, bold=True, color=DARK_BLUE)

add_rect(sl, Inches(0.4), Inches(1.65), W - Inches(0.8), Inches(1.1), LIGHT_GREY)
add_text(sl, 'Caption sinh ra:  "a dog is running in the park"',
         Inches(0.6), Inches(1.72), W - Inches(1.2), Inches(0.38),
         font_size=14, color=MID_BLUE, bold=True)
add_text(sl, 'Ground truth:      "a brown dog runs across the grass"',
         Inches(0.6), Inches(2.1), W - Inches(1.2), Inches(0.38),
         font_size=14, color=GREEN, bold=True)

bleu_rows = [
    ("BLEU-1", "1-gram", '"a", "dog", "is", "the" → 4/7 khớp', "Dễ nhất"),
    ("BLEU-2", "2-gram", '"a dog" → 1/6 bigram khớp',           "Trung bình"),
    ("BLEU-3", "3-gram", "3 từ liên tiếp phải khớp",             "Khó"),
    ("BLEU-4", "4-gram", "4 từ liên tiếp phải khớp",             "Khó nhất — CHỈ SỐ CHÍNH"),
]
col_ws_b = [Inches(1.4), Inches(1.4), Inches(5.5), Inches(3.6)]
yt_b = Inches(2.85)
add_rect(sl, Inches(0.4), yt_b, sum(col_ws_b), Inches(0.38), DARK_BLUE)
for j, h in enumerate(["Thước đo", "Loại", "Ví dụ", "Mức độ"]):
    add_text(sl, h, Inches(0.4) + sum(col_ws_b[:j]), yt_b + Inches(0.05),
             col_ws_b[j], Inches(0.28), font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
for ri, (b, ng, ex, level) in enumerate(bleu_rows):
    yy = yt_b + Inches(0.38) + ri * Inches(0.45)
    bg_b = LIGHT_GREY if ri % 2 == 0 else WHITE
    if ri == 3: bg_b = RGBColor(0xFF, 0xFF, 0xCC)
    add_rect(sl, Inches(0.4), yy, sum(col_ws_b), Inches(0.43), bg_b)
    for j, val in enumerate([b, ng, ex, level]):
        fw = True if (j == 0 or ri == 3) else False
        col_t = RED if ri == 3 else DARK_GREY
        add_text(sl, val, Inches(0.4) + sum(col_ws_b[:j]) + Inches(0.05), yy + Inches(0.07),
                 col_ws_b[j] - Inches(0.1), Inches(0.3),
                 font_size=12, bold=fw, color=col_t, align=PP_ALIGN.CENTER if j < 2 else PP_ALIGN.LEFT)

add_rect(sl, Inches(0.4), Inches(5.1), W - Inches(0.8), Inches(0.5), DARK_BLUE)
add_text(sl, "Đánh giá trên  1,000 ảnh × 5 caption/ảnh = 5,000 câu tham chiếu  →  Kết quả đáng tin cậy",
         Inches(0.6), Inches(5.18), W - Inches(1.2), Inches(0.38),
         font_size=14, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

# ─── SLIDE 18: ĐÁNH GIÁ VGG16 ───────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK_LAYOUT)
slide_bg(sl)
title_bar(sl, "KẾT QUẢ ĐÁNH GIÁ — MÔ HÌNH 1: VGG16", "Beam Width = 5  |  Test set = 1,000 ảnh")
footer(sl)

if len(imgs_vgg) > 9:
    add_img(sl, imgs_vgg[9][2], Inches(0.3), Inches(1.15), Inches(9.0), Inches(4.3))

right_x2 = Inches(9.5)
add_text(sl, "BLEU Scores:", right_x2, Inches(1.2), Inches(3.6), Inches(0.4),
         font_size=14, bold=True, color=DARK_BLUE)
vgg_bleu = [("BLEU-1", "61.6%", MID_BLUE),
            ("BLEU-2", "43.6%", MID_BLUE),
            ("BLEU-3", "30.5%", MID_BLUE),
            ("BLEU-4", "21.3%", RED)]
for j, (k, v, c) in enumerate(vgg_bleu):
    y = Inches(1.65) + j * Inches(0.55)
    add_rect(sl, right_x2, y, Inches(3.6), Inches(0.48), LIGHT_GREY if j%2==0 else WHITE)
    add_text(sl, k, right_x2 + Inches(0.1), y + Inches(0.1), Inches(1.5), Inches(0.28),
             font_size=14, bold=True, color=DARK_BLUE)
    add_text(sl, v, right_x2 + Inches(1.8), y + Inches(0.1), Inches(1.7), Inches(0.28),
             font_size=18, bold=True, color=c, align=PP_ALIGN.RIGHT)

add_rect(sl, right_x2, Inches(3.95), Inches(3.6), Inches(0.55), DARK_BLUE)
add_text(sl, "→  Baseline\nBLEU-4: 21.3%", right_x2 + Inches(0.1), Inches(4.0), Inches(3.4), Inches(0.48),
         font_size=13, bold=True, color=ORANGE)

# ─── SLIDE 19: ĐÁNH GIÁ RESNET-101 ──────────────────────────────────────────
sl = prs.slides.add_slide(BLANK_LAYOUT)
slide_bg(sl)
title_bar(sl, "KẾT QUẢ ĐÁNH GIÁ — MÔ HÌNH 2: RESNET-101", "Beam Width = 5  |  Test set = 1,000 ảnh")
footer(sl)

if len(imgs_res) > 5:
    add_img(sl, imgs_res[5][2], Inches(0.3), Inches(1.15), Inches(9.0), Inches(4.3))

right_x2 = Inches(9.5)
add_text(sl, "BLEU Scores:", right_x2, Inches(1.2), Inches(3.6), Inches(0.4),
         font_size=14, bold=True, color=DARK_BLUE)
res_bleu = [("BLEU-1", "68.1%", GREEN),
            ("BLEU-2", "49.5%", GREEN),
            ("BLEU-3", "35.0%", GREEN),
            ("BLEU-4", "24.8% 🏆", GREEN)]
for j, (k, v, c) in enumerate(res_bleu):
    y = Inches(1.65) + j * Inches(0.55)
    add_rect(sl, right_x2, y, Inches(3.6), Inches(0.48), LIGHT_GREY if j%2==0 else WHITE)
    add_text(sl, k, right_x2 + Inches(0.1), y + Inches(0.1), Inches(1.5), Inches(0.28),
             font_size=14, bold=True, color=DARK_BLUE)
    add_text(sl, v, right_x2 + Inches(1.7), y + Inches(0.1), Inches(1.8), Inches(0.28),
             font_size=16, bold=True, color=c, align=PP_ALIGN.RIGHT)

add_rect(sl, right_x2, Inches(3.95), Inches(3.6), Inches(0.55), GREEN)
add_text(sl, "→  TỐT NHẤT\n+16.4% so với VGG16", right_x2 + Inches(0.1), Inches(4.0), Inches(3.4), Inches(0.48),
         font_size=13, bold=True, color=WHITE)

# ─── SLIDE 20: ĐÁNH GIÁ CLIP ─────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK_LAYOUT)
slide_bg(sl)
title_bar(sl, "KẾT QUẢ ĐÁNH GIÁ — MÔ HÌNH 3: CLIP", "Beam Width = 5  |  Test set = 1,000 ảnh")
footer(sl)

if len(imgs_clip) > 5:
    add_img(sl, imgs_clip[5][2], Inches(0.3), Inches(1.15), Inches(9.0), Inches(4.3))

right_x2 = Inches(9.5)
add_text(sl, "BLEU Scores:", right_x2, Inches(1.2), Inches(3.6), Inches(0.4),
         font_size=14, bold=True, color=DARK_BLUE)
clip_bleu = [("BLEU-1", "66.0%", ORANGE),
             ("BLEU-2", "48.3%", ORANGE),
             ("BLEU-3", "34.4%", ORANGE),
             ("BLEU-4", "24.3%", ORANGE)]
for j, (k, v, c) in enumerate(clip_bleu):
    y = Inches(1.65) + j * Inches(0.55)
    add_rect(sl, right_x2, y, Inches(3.6), Inches(0.48), LIGHT_GREY if j%2==0 else WHITE)
    add_text(sl, k, right_x2 + Inches(0.1), y + Inches(0.1), Inches(1.5), Inches(0.28),
             font_size=14, bold=True, color=DARK_BLUE)
    add_text(sl, v, right_x2 + Inches(1.8), y + Inches(0.1), Inches(1.7), Inches(0.28),
             font_size=16, bold=True, color=c, align=PP_ALIGN.RIGHT)

add_rect(sl, right_x2, Inches(3.95), Inches(3.6), Inches(0.55), ORANGE)
add_text(sl, "→  Xếp 2\n+14.1% so với VGG16", right_x2 + Inches(0.1), Inches(4.0), Inches(3.4), Inches(0.48),
         font_size=13, bold=True, color=WHITE)

# ─── SLIDE 21: SO SÁNH 3 MÔ HÌNH ────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK_LAYOUT)
slide_bg(sl)
title_bar(sl, "SO SÁNH KẾT QUẢ 3 MÔ HÌNH", "Test set: 1,000 ảnh  |  5 caption/ảnh  |  Beam Width = 5")
footer(sl)

comp_headers = ["Mô hình", "BLEU-1", "BLEU-2", "BLEU-3", "BLEU-4", "Cải thiện vs VGG16"]
comp_rows = [
    ["VGG16 + LSTM",       "61.6%", "43.6%", "30.5%", "21.3%", "—  (Baseline)"],
    ["CLIP + LSTM",        "66.0%", "48.3%", "34.4%", "24.3%", "+14.1%"],
    ["ResNet-101 + LSTM",  "68.1%", "49.5%", "35.0%", "24.8%", "+16.4%  🏆"],
]
comp_cws = [Inches(2.8), Inches(1.3), Inches(1.3), Inches(1.3), Inches(1.7), Inches(3.0)]
yt_c = Inches(1.3)
add_rect(sl, Inches(0.4), yt_c, sum(comp_cws), Inches(0.42), DARK_BLUE)
for j, h in enumerate(comp_headers):
    add_text(sl, h, Inches(0.4) + sum(comp_cws[:j]), yt_c + Inches(0.06),
             comp_cws[j], Inches(0.3), font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
row_colors = [WHITE, RGBColor(0xFF, 0xF0, 0xE0), RGBColor(0xE8, 0xFF, 0xE8)]
for ri, row in enumerate(comp_rows):
    yy = yt_c + Inches(0.42) + ri * Inches(0.58)
    add_rect(sl, Inches(0.4), yy, sum(comp_cws), Inches(0.56), row_colors[ri])
    for j, val in enumerate(row):
        is_best = (ri == 2 and j in [1, 2, 3, 4, 5])
        col_v = GREEN if is_best else DARK_GREY
        add_text(sl, val, Inches(0.4) + sum(comp_cws[:j]) + Inches(0.05), yy + Inches(0.12),
                 comp_cws[j] - Inches(0.1), Inches(0.35),
                 font_size=13 if not is_best else 15, bold=is_best, color=col_v,
                 align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

add_rect(sl, Inches(0.4), Inches(3.15), sum(comp_cws), Inches(0.04), ORANGE)

# So sánh với nghiên cứu khác
add_text(sl, "So sánh với các công trình quốc tế (Flickr30k):", Inches(0.4), Inches(3.3),
         W - Inches(0.8), Inches(0.38), font_size=14, bold=True, color=DARK_BLUE)
pub_headers = ["Công trình", "BLEU-4"]
pub_rows = [
    ["Xu et al., 2015 — Hard Attention (VGG16)", "19.9%"],
    ["Vinyals et al., 2015 — Google (GoogLeNet)", "24.3%"],
    ["Nhóm em — VGG16 + LSTM + Attention", "21.3%  ↑ Xu et al."],
    ["Nhóm em — CLIP + LSTM + Attention", "24.3%  = Google"],
    ["Nhóm em — ResNet-101 + LSTM + Attention", "24.8%  > Google  🏆"],
]
pub_cws = [Inches(8.5), Inches(3.5)]
yt_p = Inches(3.75)
add_rect(sl, Inches(0.4), yt_p, sum(pub_cws), Inches(0.38), DARK_BLUE)
for j, h in enumerate(pub_headers):
    add_text(sl, h, Inches(0.4) + sum(pub_cws[:j]), yt_p + Inches(0.05),
             pub_cws[j], Inches(0.28), font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
for ri, row in enumerate(pub_rows):
    yy = yt_p + Inches(0.38) + ri * Inches(0.42)
    bg_p = LIGHT_GREY if ri % 2 == 0 else WHITE
    if ri >= 2: bg_p = RGBColor(0xE8, 0xFF, 0xE8)
    add_rect(sl, Inches(0.4), yy, sum(pub_cws), Inches(0.4), bg_p)
    for j, val in enumerate(row):
        is_top = (ri == 4)
        add_text(sl, val, Inches(0.4) + sum(pub_cws[:j]) + Inches(0.05), yy + Inches(0.06),
                 pub_cws[j] - Inches(0.1), Inches(0.28),
                 font_size=12, bold=is_top, color=(GREEN if is_top else DARK_GREY))

# ─── SLIDE 22: DEMO VGG16 ────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK_LAYOUT)
slide_bg(sl)
title_bar(sl, "DEMO KẾT QUẢ — MÔ HÌNH 1: VGG16", "Caption sinh ra trên ảnh thật từ tập Test")
footer(sl)

if len(imgs_vgg) > 10:
    add_img(sl, imgs_vgg[10][2], Inches(0.3), Inches(1.15), W - Inches(0.6), Inches(4.3))
if len(imgs_vgg) > 11:
    add_text(sl, "Phân tích chi tiết:", Inches(0.4), Inches(5.55), Inches(4), Inches(0.35),
             font_size=13, bold=True, color=DARK_BLUE)

# ─── SLIDE 23: DEMO RESNET-101 ───────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK_LAYOUT)
slide_bg(sl)
title_bar(sl, "DEMO KẾT QUẢ — MÔ HÌNH 2: RESNET-101", "Caption sinh ra trên ảnh thật từ tập Test")
footer(sl)

if len(imgs_res) > 6:
    add_img(sl, imgs_res[6][2], Inches(0.3), Inches(1.15), W - Inches(0.6), Inches(4.3))

# ─── SLIDE 24: DEMO CLIP ─────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK_LAYOUT)
slide_bg(sl)
title_bar(sl, "DEMO KẾT QUẢ — MÔ HÌNH 3: CLIP", "Caption sinh ra trên ảnh thật từ tập Test")
footer(sl)

if len(imgs_clip) > 6:
    add_img(sl, imgs_clip[6][2], Inches(0.3), Inches(1.15), W - Inches(0.6), Inches(4.3))

# ─── SLIDE 25: PHÂN TÍCH KẾT QUẢ ────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK_LAYOUT)
slide_bg(sl)
title_bar(sl, "PHÂN TÍCH KẾT QUẢ", "Tại sao ResNet-101 thắng? Tại sao CLIP không vượt được ResNet?")
footer(sl)

add_text(sl, "❓  Câu hỏi thú vị: CLIP mạnh hơn về lý thuyết, nhưng tại sao ResNet-101 lại thắng?",
         Inches(0.4), Inches(1.2), W - Inches(0.8), Inches(0.45),
         font_size=16, bold=True, color=DARK_BLUE)

reasons = [
    ("1", "Feature dimension",
     "ResNet-101: (49, 2048)  vs  CLIP: (49, 512)\nBahdanau Attention hoạt động tốt hơn với vector đặc trưng nhiều chiều hơn.",
     MID_BLUE),
    ("2", "Mục đích huấn luyện",
     "CLIP tối ưu cho image-text MATCHING (phân loại)\nKhông phải cho GENERATION (sinh chuỗi) → LSTM khó khai thác.",
     ORANGE),
    ("3", "Decoder không đủ mạnh",
     "LSTM 512d không tận dụng được sức mạnh semantic của CLIP\nCần Transformer decoder → hướng phát triển tiếp theo.",
     GREEN),
]
for j, (num, title_r, body, col) in enumerate(reasons):
    y = Inches(1.8) + j * Inches(1.15)
    add_rect(sl, Inches(0.4), y, Inches(0.5), Inches(0.9), col)
    add_text(sl, num, Inches(0.4), y + Inches(0.2), Inches(0.5), Inches(0.5),
             font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, Inches(0.95), y, W - Inches(1.35), Inches(0.9), LIGHT_GREY)
    add_text(sl, title_r, Inches(1.05), y + Inches(0.05), W - Inches(1.5), Inches(0.35),
             font_size=14, bold=True, color=DARK_BLUE)
    add_text(sl, body, Inches(1.05), y + Inches(0.42), W - Inches(1.5), Inches(0.42),
             font_size=12, color=DARK_GREY)

add_rect(sl, Inches(0.4), Inches(5.3), W - Inches(0.8), Inches(0.5), DARK_BLUE)
add_text(sl, "✅  Kết luận: ResNet-101 phù hợp nhất với kiến trúc LSTM + Attention trong bài toán sinh mô tả ảnh",
         Inches(0.6), Inches(5.38), W - Inches(1.2), Inches(0.38),
         font_size=14, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

# ─── SLIDE 26: ERROR ANALYSIS ─────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK_LAYOUT)
slide_bg(sl)
title_bar(sl, "PHÂN TÍCH LỖI", "Error Analysis — Trường hợp mô hình gặp khó khăn")
footer(sl)

if len(imgs_vgg) > 11:
    add_img(sl, imgs_vgg[11][2], Inches(0.3), Inches(1.15), Inches(7.0), Inches(4.2))

right_x3 = Inches(7.5)
add_text(sl, "Mô hình gặp khó với:", right_x3, Inches(1.2), Inches(5.6), Inches(0.4),
         font_size=14, bold=True, color=DARK_BLUE)
errors = [
    ("🌄", "Cảnh thiên nhiên phức tạp", "Núi, rừng, bãi biển không có người"),
    ("👥", "Nhiều người trong ảnh", "Mô tả chung chung, không cụ thể"),
    ("🎭", "Hoạt động trừu tượng", "Cảm xúc, tương tác phức tạp"),
    ("🌙", "Ánh sáng / góc chụp lạ", "Ảnh tối, ngược sáng, góc độ bất thường"),
]
for j, (icon, title_e, desc) in enumerate(errors):
    y = Inches(1.7) + j * Inches(0.75)
    add_text(sl, f"{icon}  {title_e}", right_x3, y, Inches(5.6), Inches(0.35),
             font_size=13, bold=True, color=RED)
    add_text(sl, desc, right_x3 + Inches(0.2), y + Inches(0.35), Inches(5.4), Inches(0.32),
             font_size=12, color=DARK_GREY)

add_rect(sl, right_x3, Inches(4.85), Inches(5.6), Inches(0.5), LIGHT_GREY)
add_text(sl, "→  Đây là hạn chế chung của kiến trúc LSTM\n    Transformer decoder có thể xử lý tốt hơn",
         right_x3 + Inches(0.1), Inches(4.88), Inches(5.4), Inches(0.45),
         font_size=12, color=DARK_GREY, italic=True)

# ─── SLIDE 27: KẾT LUẬN & HƯỚNG PHÁT TRIỂN ──────────────────────────────────
sl = prs.slides.add_slide(BLANK_LAYOUT)
slide_bg(sl)
title_bar(sl, "KẾT LUẬN & HƯỚNG PHÁT TRIỂN", "Những gì đã đạt được và bước tiếp theo")
footer(sl)

add_text(sl, "Những gì nhóm em đã thực hiện được:", Inches(0.4), Inches(1.2), Inches(6.2), Inches(0.4),
         font_size=15, bold=True, color=DARK_BLUE)
achieved = [
    "Xây dựng và so sánh thành công 3 mô hình Image Captioning trên Flickr30k",
    "ResNet-101 đạt BLEU-4 = 24.8%  —  tốt nhất trong 3 mô hình",
    "Kết quả vượt Xu et al. 2015 (19.9%)  và  tương đương Google 2015 (24.3%)",
    "Phát hiện: CLIP không vượt ResNet với LSTM decoder — mở ra câu hỏi nghiên cứu",
    "Xây dựng Streamlit demo app — upload ảnh, sinh caption real-time",
]
for j, item in enumerate(achieved):
    y = Inches(1.65) + j * Inches(0.48)
    add_rect(sl, Inches(0.4), y, Inches(0.4), Inches(0.38), GREEN)
    add_text(sl, "✓", Inches(0.4), y + Inches(0.05), Inches(0.4), Inches(0.3),
             font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, item, Inches(0.9), y + Inches(0.05), Inches(5.5), Inches(0.35),
             font_size=13, color=DARK_GREY)

add_text(sl, "Hướng phát triển tiếp theo:", Inches(6.8), Inches(1.2), Inches(6.3), Inches(0.4),
         font_size=15, bold=True, color=DARK_BLUE)
future = [
    ("Transformer Decoder", "Thay LSTM → CLIP phát huy toàn bộ sức mạnh"),
    ("Fine-tune Encoder", "Không frozen encoder → BLEU tăng thêm ~2%"),
    ("Dataset lớn hơn", "COCO 330K ảnh → BLEU-4 có thể >30%"),
    ("CIDEr Optimization", "Loss tốt hơn Cross-entropy cho captioning"),
    ("Scheduled Sampling", "Giảm exposure bias trong training"),
]
for j, (title_f, desc) in enumerate(future):
    y = Inches(1.65) + j * Inches(0.65)
    add_rect(sl, Inches(6.8), y, Inches(2.0), Inches(0.58), DARK_BLUE)
    add_text(sl, title_f, Inches(6.85), y + Inches(0.1), Inches(1.9), Inches(0.38),
             font_size=12, bold=True, color=WHITE)
    add_text(sl, desc, Inches(8.85), y + Inches(0.1), Inches(4.2), Inches(0.38),
             font_size=12, color=DARK_GREY)

# ─── SLIDE 28: CẢM ƠN ───────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(sl, 0, 0, W, H, DARK_BLUE)
add_rect(sl, 0, Inches(2.5), W, Inches(0.06), ORANGE)
add_rect(sl, 0, Inches(5.2), W, Inches(0.06), ORANGE)

add_text(sl, "XIN TRÂN TRỌNG CẢM ƠN",
         Inches(0.5), Inches(1.2), W - Inches(1), Inches(0.8),
         font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Hội đồng và thầy Hồ Nhựt Minh đã lắng nghe",
         Inches(0.5), Inches(1.95), W - Inches(1), Inches(0.5),
         font_size=18, color=ORANGE, align=PP_ALIGN.CENTER, italic=True)

add_text(sl, "Nhóm xin sẵn sàng trả lời câu hỏi",
         Inches(0.5), Inches(2.7), W - Inches(1), Inches(0.6),
         font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(sl, "Lê Văn Cảnh    |    Nguyễn Đức Trường Giang",
         Inches(0.5), Inches(3.5), W - Inches(1), Inches(0.5),
         font_size=20, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "DT2307L  —  Aptech",
         Inches(0.5), Inches(4.0), W - Inches(1), Inches(0.4),
         font_size=16, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

# Tóm tắt kết quả cuối
add_text(sl, "ResNet-101: BLEU-4 = 24.8% 🏆     CLIP: BLEU-4 = 24.3%     VGG16: BLEU-4 = 21.3%",
         Inches(0.5), Inches(5.4), W - Inches(1), Inches(0.45),
         font_size=15, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

# ─── Lưu file ────────────────────────────────────────────────────────────────
output_path = 'BAO_VE_DO_AN.pptx'
prs.save(output_path)
print(f"\n✅  Đã tạo: {output_path}")
print(f"   Số slides: {len(prs.slides)}")
print(f"   Ảnh nhúng từ VGG16 notebook:   {len(imgs_vgg)}")
print(f"   Ảnh nhúng từ ResNet101 notebook:{len(imgs_res)}")
print(f"   Ảnh nhúng từ CLIP notebook:     {len(imgs_clip)}")
print(f"   Tổng ảnh:                       {len(imgs_vgg)+len(imgs_res)+len(imgs_clip)}")
